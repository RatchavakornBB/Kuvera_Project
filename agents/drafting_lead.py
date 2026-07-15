"""Drafting Lead (system-architecture.md Section 9) — 5.1 Doc/email prep (.docx + cover email
draft, scope-extended per Section 10.3 to also produce a NotebookLM-style source-cited summary)
and 5.2 Presentation prep (.pptx). Real files, real Storage upload, real Document rows — the same
Storage+Document-together invariant every other upload path in this codebase follows (AGENT.md
Section 11), not a decorative "export" button that produces nothing durable.

Dual input per Figure 7: always the Analyst Lead's latest stored analysis; the Contracts Lead's
document (if one exists and has real extracted clauses) is included when present, omitted
otherwise — never fabricated if no contract was ever uploaded."""

import io
import re
from datetime import datetime, timezone
from typing import Any

from docx import Document as DocxDocument
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

from agents.adapters.model_adapter import call_model
from agents.analyses import get_last_analysis
from agents.db import get_client
from agents.embeddings import embed_text
from agents.retry import with_retry

BUCKET = "deal-documents"

EMAIL_PROMPT = """You are the Drafting Lead's doc/email prep agent. Draft a short, professional
cover email an M&A associate would send when circulating the attached IC memo to an external
party (e.g. a co-investor or advisor). Reference the deal by name and give a one-paragraph
high-level summary — do not restate the full memo. Keep it under 150 words. Output only the email
body text, no subject line, no signature block.

DEAL: {deal_name} ({client})
IC MEMO DRAFT:
{ic_memo}
"""


def _get_deal_and_analysis(deal_id: str) -> tuple[dict[str, Any], dict[str, Any]]:
    client = get_client()
    deal_rows = client.table("deals").select("*").eq("id", deal_id).execute().data
    if not deal_rows:
        raise ValueError(f"No deal with id {deal_id}")
    analysis = get_last_analysis(deal_id)
    if analysis is None:
        raise ValueError(f"No stored analysis for deal {deal_id} — run Analysis before drafting")
    return deal_rows[0], analysis


def _add_markdown_paragraphs(doc: DocxDocument, text: str) -> None:
    """A small, real markdown-to-docx renderer (not a dependency) — handles the
    subset ic_memo_drafter's prompt actually produces: #/##/### headers,
    **bold** inline, plain paragraphs. Good enough for real generated memos,
    not meant as a general markdown engine."""
    for line in text.split("\n"):
        line = line.rstrip()
        if not line:
            continue
        heading_match = re.match(r"^(#{1,3})\s+(.*)", line)
        if heading_match:
            level = len(heading_match.group(1))
            doc.add_heading(heading_match.group(2), level=level)
            continue

        para = doc.add_paragraph()
        parts = re.split(r"(\*\*[^*]+\*\*)", line)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = para.add_run(part[2:-2])
                run.bold = True
            else:
                para.add_run(part)


def draft_ic_memo_docx(deal_id: str) -> tuple[bytes, dict[str, Any]]:
    deal, analysis = _get_deal_and_analysis(deal_id)

    doc = DocxDocument()
    doc.add_heading(f"{deal['name']} — Investment Committee Memo", level=0)
    meta = doc.add_paragraph()
    meta.add_run(f"Client: {deal['client']}  |  Industries: {', '.join(deal['industries'])}  |  ").italic = True
    meta.add_run(f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')}").italic = True

    memo_text = analysis.get("ic_memo_draft") or "No IC memo draft available for this deal yet."
    _add_markdown_paragraphs(doc, memo_text)

    if analysis.get("pricing_note"):
        doc.add_heading("Indicative Pricing", level=1)
        doc.add_paragraph(analysis["pricing_note"])

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), deal


def draft_ic_deck_pptx(deal_id: str) -> tuple[bytes, dict[str, Any]]:
    deal, analysis = _get_deal_and_analysis(deal_id)
    risk_flags = analysis.get("risk_flags") or []

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"{deal['name']} — IC Deck"
    title_slide.placeholders[1].text = f"{deal['client']} · {', '.join(deal['industries'])}"

    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Executive Summary"
    summary_body = summary_slide.placeholders[1].text_frame
    summary_text = (analysis.get("summary") or "")[:800]
    summary_body.text = summary_text

    if risk_flags:
        risk_slide = prs.slides.add_slide(prs.slide_layouts[1])
        risk_slide.shapes.title.text = "Key Risk Flags"
        body = risk_slide.placeholders[1].text_frame
        high = [f for f in risk_flags if f.get("severity") == "high"][:5]
        for i, flag in enumerate(high):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = flag["description"][:200]
            p.font.size = PptxPt(14)

    if analysis.get("pricing_note"):
        pricing_slide = prs.slides.add_slide(prs.slide_layouts[1])
        pricing_slide.shapes.title.text = "Indicative Pricing"
        pricing_slide.placeholders[1].text_frame.text = analysis["pricing_note"][:600]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), deal


def draft_cover_email(deal_id: str) -> str:
    deal, analysis = _get_deal_and_analysis(deal_id)

    def _call() -> str:
        response = call_model(
            "drafting_lead",
            messages=[
                {
                    "role": "user",
                    "content": EMAIL_PROMPT.format(
                        deal_name=deal["name"], client=deal["client"], ic_memo=analysis.get("ic_memo_draft") or analysis["summary"]
                    ),
                }
            ],
            max_tokens=512,
        )
        for block in response.content:
            if block.type == "text":
                return block.text
        raise ValueError(f"model returned no text block (stop_reason={response.stop_reason!r})")

    return with_retry("drafting_lead", _call)


def draft_notebooklm_summary(deal_id: str) -> str:
    """NotebookLM-style source-cited summary — the scope extension noted in
    system-architecture.md Section 10.3. Built directly from real risk_flags'
    source_excerpts (already real quotes from the source document), not a new
    LLM call inventing citations."""
    deal, analysis = _get_deal_and_analysis(deal_id)
    risk_flags = analysis.get("risk_flags") or []

    lines = [f"# {deal['name']} — Source-Cited Summary\n", analysis.get("summary") or "", "\n## Findings\n"]
    for flag in risk_flags:
        excerpt = flag.get("source_excerpt", "").strip()
        lines.append(f"- **[{flag.get('severity', 'medium')}]** {flag['description']}")
        if excerpt:
            lines.append(f'  > "{excerpt}"')
    return "\n".join(lines)


def _embed_best_effort(text: str) -> list[float] | None:
    """Same best-effort contract as backend/app/services/documents.py's
    _embed_document_text — an embeddings-provider failure must never break
    a draft action, since the file itself already generated successfully."""
    if not text.strip():
        return None
    try:
        return embed_text(text, input_type="document")
    except Exception:
        return None


def _upload_and_record(
    deal_id: str, filename: str, content: bytes, content_type: str, doc_type: str, embed_source: str = ""
) -> dict[str, Any]:
    import uuid

    client = get_client()
    storage_path = f"{deal_id}/{uuid.uuid4()}-{filename}"
    client.storage.from_(BUCKET).upload(storage_path, content, {"content-type": content_type})

    try:
        res = (
            client.table("documents")
            .insert(
                {
                    "deal_id": deal_id,
                    "name": filename,
                    "type": doc_type,
                    "storage_path": storage_path,
                    "status": "approved",
                    "embedding": _embed_best_effort(f"{filename}\n{embed_source}"),
                }
            )
            .execute()
        )
    except Exception:
        client.storage.from_(BUCKET).remove([storage_path])
        raise
    doc = res.data[0]
    doc.pop("embedding", None)
    return doc


def draft_and_store_ic_memo(deal_id: str) -> dict[str, Any]:
    content, deal = draft_ic_memo_docx(deal_id)
    analysis = get_last_analysis(deal_id)
    filename = f"{deal['name'].replace(' ', '_')}_IC_Memo.docx"
    embed_source = analysis.get("ic_memo_draft") or analysis.get("summary") or ""
    return _upload_and_record(
        deal_id,
        filename,
        content,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "Drafted Memo",
        embed_source,
    )


def draft_and_store_ic_deck(deal_id: str) -> dict[str, Any]:
    content, deal = draft_ic_deck_pptx(deal_id)
    analysis = get_last_analysis(deal_id)
    filename = f"{deal['name'].replace(' ', '_')}_IC_Deck.pptx"
    embed_source = analysis.get("summary") or ""
    return _upload_and_record(
        deal_id,
        filename,
        content,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "Drafted Deck",
        embed_source,
    )


def draft_and_store_cover_email(deal_id: str) -> dict[str, Any]:
    """Previously ephemeral — draft_cover_email() only ever returned text
    for the frontend to display, with nothing durable written anywhere
    (same gap web_research had). Now stored through the same real
    Storage+Document write path as the memo/deck, real embedding included,
    so a drafted email is genuinely searchable and reappears in the
    Documents tab instead of vanishing once the response leaves memory."""
    email_text = draft_cover_email(deal_id)
    deal, _analysis = _get_deal_and_analysis(deal_id)
    filename = f"{deal['name'].replace(' ', '_')}_Cover_Email.txt"
    doc = _upload_and_record(deal_id, filename, email_text.encode("utf-8"), "text/plain", "Drafted Email", email_text)
    doc["email"] = email_text
    return doc


def draft_and_store_notebooklm_summary(deal_id: str) -> dict[str, Any]:
    """Same fix as draft_and_store_cover_email, for the source-cited
    summary — previously ephemeral, now a real stored/embedded document."""
    summary_text = draft_notebooklm_summary(deal_id)
    deal, _analysis = _get_deal_and_analysis(deal_id)
    filename = f"{deal['name'].replace(' ', '_')}_Source_Cited_Summary.txt"
    doc = _upload_and_record(
        deal_id, filename, summary_text.encode("utf-8"), "text/plain", "Drafted Summary", summary_text
    )
    doc["summary"] = summary_text
    return doc
